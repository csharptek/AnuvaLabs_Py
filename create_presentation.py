#!/usr/bin/env python3
"""
Script to create a PowerPoint presentation for the Security Testing API architecture
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import os

def create_presentation():
    """Create a comprehensive PowerPoint presentation"""
    
    # Create presentation
    prs = Presentation()
    
    # Slide 1: Title Slide
    slide_layout = prs.slide_layouts[0]  # Title slide layout
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "Security Testing API"
    subtitle.text = "Architecture & Design Documentation\n\nFastAPI-based Security Testing Platform\nwith JWT Authentication"
    
    # Slide 2: Overview
    slide_layout = prs.slide_layouts[1]  # Title and content layout
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Project Overview"
    content.text = """• FastAPI-based web application for security testing
• JWT authentication with access and refresh tokens
• ZIP file upload and vulnerability scanning
• Mock security analysis with structured reporting
• RESTful API design with comprehensive documentation
• Modular architecture with clear separation of concerns"""
    
    # Slide 3: System Architecture
    slide_layout = prs.slide_layouts[5]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Add title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "System Architecture"
    title_frame.paragraphs[0].font.size = Pt(28)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Add architecture layers
    layers = [
        {"name": "Client Layer", "desc": "Client Apps & Swagger UI", "color": RGBColor(225, 245, 254), "y": 1.5},
        {"name": "API Layer", "desc": "FastAPI Application with Routers", "color": RGBColor(243, 229, 245), "y": 2.5},
        {"name": "Service Layer", "desc": "Authentication & Security Services", "color": RGBColor(232, 245, 232), "y": 3.5},
        {"name": "Data Layer", "desc": "Models, JWT Helper & Mock Database", "color": RGBColor(252, 228, 236), "y": 4.5}
    ]
    
    for layer in layers:
        # Layer box
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(1), Inches(layer["y"]), Inches(8), Inches(0.8)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = layer["color"]
        shape.line.color.rgb = RGBColor(0, 0, 0)
        
        # Layer text
        text_frame = shape.text_frame
        text_frame.text = f"{layer['name']}: {layer['desc']}"
        text_frame.paragraphs[0].font.size = Pt(16)
        text_frame.paragraphs[0].font.bold = True
        text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Slide 4: API Endpoints
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "API Endpoints"
    content.text = """Authentication Endpoints:
• POST /login - User authentication with username/password
• POST /refresh-token - Refresh expired JWT tokens
• GET /me - Get current authenticated user information

Security Testing Endpoints:
• POST /api/v1/security-testing - Upload ZIP file for scanning
• GET / - Root endpoint with API information
• GET /docs - Interactive Swagger UI documentation

Response Formats:
• JSON responses with structured error handling
• HTTP status codes (200, 400, 401, 500)
• Comprehensive API documentation"""
    
    # Slide 5: Authentication Flow
    slide_layout = prs.slide_layouts[5]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Add title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "Authentication Flow"
    title_frame.paragraphs[0].font.size = Pt(28)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Flow steps
    steps = [
        "1. Client sends POST /login with credentials",
        "2. Auth Router validates request format",
        "3. Auth Service authenticates user against database",
        "4. JWT Helper creates access & refresh tokens",
        "5. Tokens returned to client (60min + 7days)",
        "6. Client uses Bearer token for protected endpoints",
        "7. Token validation on each protected request"
    ]
    
    y_pos = 1.8
    for step in steps:
        text_box = slide.shapes.add_textbox(Inches(1), Inches(y_pos), Inches(8), Inches(0.4))
        text_frame = text_box.text_frame
        text_frame.text = step
        text_frame.paragraphs[0].font.size = Pt(14)
        y_pos += 0.5
    
    # Slide 6: Security Testing Flow
    slide_layout = prs.slide_layouts[5]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Add title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "Security Testing Flow"
    title_frame.paragraphs[0].font.size = Pt(28)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Flow steps
    steps = [
        "1. Client uploads ZIP file via multipart/form-data",
        "2. Security Router validates file extension (.zip)",
        "3. Security Service creates temporary directory",
        "4. ZIP file extracted and contents analyzed",
        "5. Each file processed for vulnerability patterns",
        "6. Mock issues generated based on file type",
        "7. Scan results compiled and returned",
        "8. Temporary files cleaned up securely"
    ]
    
    y_pos = 1.8
    for step in steps:
        text_box = slide.shapes.add_textbox(Inches(1), Inches(y_pos), Inches(8), Inches(0.4))
        text_frame = text_box.text_frame
        text_frame.text = step
        text_frame.paragraphs[0].font.size = Pt(14)
        y_pos += 0.4
    
    # Slide 7: Data Models
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Data Models"
    content.text = """Authentication Models:
• UserBase: username, email, role
• UserInDB: extends UserBase + id, password
• User: extends UserBase + id (no password)
• UserLogin: username, password
• Token: access_token, refresh_token, token_type
• TokenPayload: sub, exp, type

Security Testing Models:
• VulnerabilityItem: file, issues[]
• ScanResponse: status, file_count, vulnerabilities[]
• ErrorResponse: status, message, details

All models use Pydantic for validation and serialization"""
    
    # Slide 8: Component Dependencies
    slide_layout = prs.slide_layouts[5]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Add title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "Component Dependencies"
    title_frame.paragraphs[0].font.size = Pt(28)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Components
    components = [
        "main.py (FastAPI App) → Routers",
        "auth/routes.py → auth/service.py",
        "app/api/v1/security.py → app/services/security_service.py",
        "auth/service.py → auth/jwt_helper.py + auth/models.py",
        "security_service.py → app/models/security.py",
        "jwt_helper.py → config.py",
        "All components → Pydantic models for validation"
    ]
    
    y_pos = 2
    for comp in components:
        text_box = slide.shapes.add_textbox(Inches(1), Inches(y_pos), Inches(8), Inches(0.5))
        text_frame = text_box.text_frame
        text_frame.text = comp
        text_frame.paragraphs[0].font.size = Pt(14)
        y_pos += 0.6
    
    # Slide 9: Security Features
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Security Features"
    content.text = """Current Implementation:
• JWT-based authentication with token expiration
• Separate access (60min) and refresh (7days) tokens
• Token type validation (access vs refresh)
• CORS middleware for cross-origin requests
• File type validation for uploads
• Secure temporary file handling with cleanup
• Structured error responses with appropriate HTTP codes

Production Recommendations:
• Password hashing (bcrypt/scrypt)
• Rate limiting for authentication endpoints
• Real database integration (PostgreSQL/MongoDB)
• Environment-based configuration
• Comprehensive audit logging
• Real security scanning tools integration"""
    
    # Slide 10: File Structure
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Project File Structure"
    content.text = """AnuvaLabs_Py/
├── main.py                     # FastAPI application entry point
├── config.py                   # Configuration settings
├── test_api.py                 # API testing script
├── auth/                       # Authentication module
│   ├── routes.py              # Authentication endpoints
│   ├── service.py             # Authentication business logic
│   ├── models.py              # Authentication data models
│   └── jwt_helper.py          # JWT token utilities
└── app/                       # Application module
    ├── api/v1/security.py     # Security testing endpoints
    ├── models/security.py     # Security testing models
    └── services/security_service.py # Security testing logic"""
    
    # Slide 11: Usage Examples
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Usage Examples"
    content.text = """Login Request:
POST /login
{
  "username": "admin",
  "password": "12345"
}

Response:
{
  "access_token": "eyJhbGciOiJIUzI1NiIs...",
  "refresh_token": "eyJhbGciOiJIUzI1NiIs...",
  "token_type": "bearer"
}

Security Testing:
POST /api/v1/security-testing
Content-Type: multipart/form-data
file: test.zip

Response:
{
  "status": "success",
  "file_count": 3,
  "vulnerabilities": [...]
}"""
    
    # Slide 12: Development Setup
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Development Setup"
    content.text = """Prerequisites:
• Python 3.8+
• FastAPI
• python-jose[cryptography]
• python-multipart
• uvicorn

Installation:
pip install fastapi uvicorn python-jose[cryptography] python-multipart

Run Application:
uvicorn main:app --host 0.0.0.0 --port 8000 --reload

Access Points:
• API: http://localhost:8000
• Documentation: http://localhost:8000/docs
• Testing: python test_api.py"""
    
    # Slide 13: Future Enhancements
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Future Enhancements"
    content.text = """Database Integration:
• PostgreSQL/MongoDB for persistent storage
• User management with registration/password reset

Security Improvements:
• Real security scanning tools (SAST/DAST)
• Enhanced file type validation
• Rate limiting and DDoS protection

Features:
• Role-based access control
• Audit logging and monitoring
• Background job processing for large files
• Support for additional file formats
• Caching with Redis
• Docker containerization
• Health checks and metrics collection"""
    
    # Save presentation
    prs.save('Security_Testing_API_Architecture.pptx')
    print("✅ PowerPoint presentation saved as 'Security_Testing_API_Architecture.pptx'")

if __name__ == "__main__":
    print("📊 Creating PowerPoint presentation...")
    create_presentation()
    print("🎉 Presentation created successfully!")

